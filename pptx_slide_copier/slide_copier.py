"""Slide copying utilities."""

from __future__ import annotations

from copy import deepcopy
from io import BytesIO

from pptx import Presentation
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.parts.slide import SlideLayoutPart, SlideMasterPart
from pptx.slide import Slide

# Relationship types that represent structural links and should be
# skipped when generically copying part relationships.
_STRUCTURAL_REL_TYPES = frozenset({RT.SLIDE_MASTER, RT.SLIDE_LAYOUT, RT.THEME})

# Namespace URI used for r:embed / r:link / r:id attributes
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


class SlideCopier:
    """Handles copying slides between presentations."""

    @staticmethod
    def copy_layouts(source_prs: Presentation, target_prs: Presentation) -> dict:
        """Copy all masters/layouts/themes from source to target at once.

        Call this once before copying slides so that the target ends up
        with only the original target themes plus the source themes.

        Args:
            source_prs: Source presentation
            target_prs: Target presentation

        Returns:
            layout_map: dict mapping {source layout name: target SlideLayout}
        """
        cache: dict = {}
        layout_map: dict[str, object] = {}

        for source_master in source_prs.slide_masters:
            source_master_part = source_master.part
            matching_target_master = SlideCopier._find_matching_master(
                source_master_part, target_prs,
            )

            if matching_target_master is not None:
                # Same theme — look up existing layouts by name, copy only missing ones
                target_master_part = matching_target_master.part
                cache[id(source_master_part)] = target_master_part
                existing = {sl.name: sl for sl in matching_target_master.slide_layouts}
                for layout in source_master.slide_layouts:
                    if layout.name in existing:
                        layout_map[layout.name] = existing[layout.name]
                    else:
                        # Layout not in target — copy it
                        source_layout_part = layout.part
                        target_layout_part = SlideCopier._copy_slide_layout_part(
                            source_layout_part, target_prs, cache,
                        )
                        cache[id(source_layout_part)] = target_layout_part
                        layout_map[layout.name] = target_layout_part.slide_layout
            else:
                # Different theme — copy everything as before
                target_master_part = SlideCopier._get_or_copy_slide_master(
                    source_master_part, target_prs, cache,
                )
                for layout in source_master.slide_layouts:
                    source_layout_part = layout.part
                    cache_key = id(source_layout_part)
                    if cache_key not in cache:
                        target_layout_part = SlideCopier._copy_slide_layout_part(
                            source_layout_part, target_prs, cache,
                        )
                        cache[cache_key] = target_layout_part
                    else:
                        target_layout_part = cache[cache_key]
                    layout_map[layout.name] = target_layout_part.slide_layout

        return layout_map

    @staticmethod
    def copy_slides(
        source_prs: Presentation,
        target_prs: Presentation,
        slide_indices=None,
        target_slide_index: int | None = None,
    ):
        """Copy multiple slides from source to target.

        Copies all source layouts/masters/themes once upfront via
        ``copy_layouts()``, then copies each slide referencing the
        pre-copied layouts by name.

        Args:
            source_prs: Source presentation
            target_prs: Target presentation
            slide_indices: List of 0-based slide indices to copy.
                           If None, all slides are copied.
            target_slide_index: Optional 0-based index at which to insert the
                           slides in the target presentation.  Slides are
                           inserted sequentially starting from this position.
                           If None, slides are appended at the end.

        Returns:
            List of newly created slides in target presentation
        """
        if slide_indices is None:
            slide_indices = list(range(len(source_prs.slides)))

        layout_map = SlideCopier.copy_layouts(source_prs, target_prs)
        slides = []
        for i, idx in enumerate(slide_indices):
            insert_at = None
            if target_slide_index is not None:
                insert_at = target_slide_index + i
            slide = SlideCopier.copy_slide(
                source_prs, idx, target_prs,
                _layout_map=layout_map,
                target_slide_index=insert_at,
            )
            slides.append(slide)
        return slides

    @staticmethod
    def copy_slide(
        source_prs: Presentation,
        source_slide_index: int,
        target_prs: Presentation,
        _layout_map: dict | None = None,
        target_slide_index: int | None = None,
    ) -> Slide:
        """Copy a slide from source presentation to target presentation.

        The source slide's layout, master, and theme are copied into the
        target at the OPC-package level so the appearance is faithfully
        preserved even when the two presentations use different templates.

        Args:
            source_prs: Source presentation
            source_slide_index: Index of slide to copy (0-based)
            target_prs: Target presentation
            _layout_map: Optional {layout_name: SlideLayout} dict returned by
                         ``copy_layouts()``.  When provided the pre-copied
                         layout is looked up by name.  When *None* the layout
                         is copied on demand (backward-compatible behaviour).
            target_slide_index: Optional 0-based index at which to insert the
                         slide in the target presentation.  When *None* (the
                         default) the slide is appended at the end.

        Returns:
            The newly created slide in target presentation
        """
        source_slide = source_prs.slides[source_slide_index]

        # Copy slide size only when target has no existing slides
        if len(target_prs.slides) == 0:
            SlideCopier._copy_slide_size(source_prs, target_prs)

        # Resolve the target layout
        if _layout_map is not None:
            source_layout_name = source_slide.slide_layout.name
            target_layout = _layout_map[source_layout_name]
        else:
            # Backward-compatible on-demand copy
            cache: dict = {}
            target_layout_part = SlideCopier._get_or_copy_slide_layout(
                source_slide, target_prs, cache,
            )
            target_layout = target_layout_part.slide_layout

        # Create new slide with the copied layout
        dest_slide = target_prs.slides.add_slide(target_layout)

        # Remove auto-generated placeholder shapes from the layout
        # to avoid duplicates when we deepcopy the source shapes below.
        spTree = dest_slide.shapes._spTree
        for sp in list(spTree.iterchildren(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}sp',
        )):
            spTree.remove(sp)

        # Copy all shapes using deepcopy at XML level
        for shape in source_slide.shapes:
            try:
                new_element = deepcopy(shape.element)
                dest_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")
            except Exception:
                continue

        # Copy all non-structural relationships (images, charts, media, etc.)
        # and remap rIds in the copied XML so references stay valid.
        rid_mapping = SlideCopier._copy_part_rels(
            source_slide.part, dest_slide.part, target_prs.part.package,
        )
        if rid_mapping:
            SlideCopier._remap_rids(dest_slide.shapes._spTree, rid_mapping)

        # Move slide to the requested position if target_slide_index is given
        if target_slide_index is not None:
            SlideCopier._move_slide_to_index(target_prs, target_slide_index)

        return dest_slide

    # ------------------------------------------------------------------
    # Layout / Master / Theme copying
    # ------------------------------------------------------------------

    @staticmethod
    def _get_or_copy_slide_layout(source_slide, target_prs, cache):
        """Return a SlideLayoutPart in target_prs that mirrors the source slide's layout.

        If the layout was already copied (present in cache), the cached part
        is returned.
        """
        source_layout_part = source_slide.part.part_related_by(RT.SLIDE_LAYOUT)
        cache_key = id(source_layout_part)
        if cache_key in cache:
            return cache[cache_key]

        target_layout_part = SlideCopier._copy_slide_layout_part(
            source_layout_part, target_prs, cache,
        )
        cache[cache_key] = target_layout_part
        return target_layout_part

    @staticmethod
    def _find_matching_master(source_master_part, target_prs):
        """Find a target master whose theme blob matches the source master's."""
        try:
            source_theme_blob = source_master_part.part_related_by(RT.THEME).blob
        except KeyError:
            return None

        for target_master in target_prs.slide_masters:
            try:
                target_theme_blob = target_master.part.part_related_by(RT.THEME).blob
            except KeyError:
                continue
            if source_theme_blob == target_theme_blob:
                return target_master

        return None

    @staticmethod
    def _copy_slide_layout_part(source_layout_part, target_prs, cache):
        """Deep-copy a SlideLayoutPart into target_prs."""
        package = target_prs.part.package

        # 1. Ensure the parent master exists in the target
        source_master_part = source_layout_part.part_related_by(RT.SLIDE_MASTER)
        target_master_part = SlideCopier._get_or_copy_slide_master(
            source_master_part, target_prs, cache,
        )

        # 2. Deep-copy the layout XML
        new_layout_element = deepcopy(source_layout_part._element)

        # 3. Create the new SlideLayoutPart
        partname = package.next_partname("/ppt/slideLayouts/slideLayout%d.xml")
        target_layout_part = SlideLayoutPart(
            partname, CT.PML_SLIDE_LAYOUT, package, new_layout_element,
        )

        # 4. Layout → Master relationship
        target_layout_part.relate_to(target_master_part, RT.SLIDE_MASTER)

        # 5. Master → Layout relationship + sldLayoutIdLst entry
        #    Register first so that the layout part is reachable via
        #    iter_parts() and image deduplication works in step 6.
        rId = target_master_part.relate_to(target_layout_part, RT.SLIDE_LAYOUT)
        sld_layout_id_lst = target_master_part._element.get_or_add_sldLayoutIdLst()
        sld_layout_id = sld_layout_id_lst._add_sldLayoutId(rId=rId)
        sld_layout_id.set("id", str(SlideCopier._next_unique_id(target_prs)))

        # 6. Copy non-structural relationships (images, etc.) and remap rIds
        rid_mapping = SlideCopier._copy_part_rels(
            source_layout_part, target_layout_part, package,
        )
        if rid_mapping:
            SlideCopier._remap_rids(new_layout_element, rid_mapping)

        return target_layout_part

    @staticmethod
    def _get_or_copy_slide_master(source_master_part, target_prs, cache):
        """Return a SlideMasterPart in target_prs that mirrors source_master_part.

        Uses cache to avoid duplicating masters.  When the source master's
        theme already exists in the target presentation the existing master
        is reused instead of creating a duplicate.
        """
        cache_key = id(source_master_part)
        if cache_key in cache:
            return cache[cache_key]

        # Check if the target already has a master with the same theme
        matching_master = SlideCopier._find_matching_master(
            source_master_part, target_prs,
        )
        if matching_master is not None:
            target_master_part = matching_master.part
            cache[cache_key] = target_master_part
            return target_master_part

        target_master_part = SlideCopier._copy_slide_master_part(
            source_master_part, target_prs, cache,
        )
        cache[cache_key] = target_master_part
        return target_master_part

    @staticmethod
    def _copy_slide_master_part(source_master_part, target_prs, cache):
        """Deep-copy a SlideMasterPart into target_prs."""
        package = target_prs.part.package

        # 1. Deep-copy the master XML
        new_master_element = deepcopy(source_master_part._element)

        # 2. Remove sldLayoutIdLst — it will be rebuilt as layouts are copied
        existing_id_lst = new_master_element.sldLayoutIdLst
        if existing_id_lst is not None:
            new_master_element.remove(existing_id_lst)

        # 3. Create the new SlideMasterPart
        partname = package.next_partname("/ppt/slideMasters/slideMaster%d.xml")
        target_master_part = SlideMasterPart(
            partname, CT.PML_SLIDE_MASTER, package, new_master_element,
        )

        # 4. Copy theme
        SlideCopier._copy_theme_part(
            source_master_part, target_master_part, target_prs, cache,
        )

        # 5. Register master in presentation's sldMasterIdLst first so
        #    that the part is reachable via iter_parts() and image
        #    deduplication works correctly in step 6.
        prs_part = target_prs.part
        rId = prs_part.relate_to(target_master_part, RT.SLIDE_MASTER)
        sld_master_id_lst = prs_part._element.get_or_add_sldMasterIdLst()
        sld_master_id = sld_master_id_lst._add_sldMasterId(rId=rId)
        sld_master_id.set("id", str(SlideCopier._next_unique_id(target_prs)))

        # 6. Copy non-structural relationships (images, etc.) and remap rIds
        rid_mapping = SlideCopier._copy_part_rels(
            source_master_part, target_master_part, package,
        )
        if rid_mapping:
            SlideCopier._remap_rids(new_master_element, rid_mapping)

        return target_master_part

    @staticmethod
    def _copy_theme_part(source_master_part, target_master_part, target_prs, cache):
        """Copy the theme part from source master to target master."""
        try:
            source_theme_part = source_master_part.part_related_by(RT.THEME)
        except KeyError:
            return

        cache_key = id(source_theme_part)
        if cache_key in cache:
            target_master_part.relate_to(cache[cache_key], RT.THEME)
            return

        package = target_prs.part.package
        partname = package.next_partname("/ppt/theme/theme%d.xml")

        # Theme parts are loaded as plain Part (blob-based) by python-pptx
        # because OFC_THEME is not in the PartFactory registry.
        target_theme_part = Part(
            partname, CT.OFC_THEME, package, blob=source_theme_part.blob,
        )

        target_master_part.relate_to(target_theme_part, RT.THEME)

        # Copy theme relationships (e.g. background images referenced via
        # r:embed inside the theme XML).  The generic _copy_part_rels cannot
        # be used here because theme is a plain blob-based Part without
        # get_or_add_image_part.
        rid_mapping: dict[str, str] = {}
        for rId, rel in source_theme_part.rels.items():
            if rel.is_external:
                new_rId = target_theme_part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True,
                )
            elif rel.reltype == RT.IMAGE:
                src_img = rel.target_part
                new_partname = package.next_partname(
                    _partname_to_template(src_img.partname),
                )
                new_img = Part(
                    new_partname, src_img.content_type,
                    package, blob=src_img.blob,
                )
                new_rId = target_theme_part.relate_to(new_img, rel.reltype)
            else:
                src_target = rel.target_part
                new_partname = package.next_partname(
                    _partname_to_template(src_target.partname),
                )
                new_part = Part(
                    new_partname, src_target.content_type,
                    package, blob=src_target.blob,
                )
                new_rId = target_theme_part.relate_to(new_part, rel.reltype)

            if rId != new_rId:
                rid_mapping[rId] = new_rId

        if rid_mapping:
            theme_xml = target_theme_part.blob.decode("utf-8")
            for old_rid, new_rid in rid_mapping.items():
                theme_xml = theme_xml.replace(
                    f'r:embed="{old_rid}"', f'r:embed="{new_rid}"',
                )
                theme_xml = theme_xml.replace(
                    f'r:link="{old_rid}"', f'r:link="{new_rid}"',
                )
            target_theme_part._blob = theme_xml.encode("utf-8")

        cache[cache_key] = target_theme_part

    # ------------------------------------------------------------------
    # Generic relationship copying & rId remapping
    # ------------------------------------------------------------------

    @staticmethod
    def _copy_part_rels(source_part, target_part, target_package):
        """Copy non-structural relationships from source_part to target_part.

        Returns a dict mapping old rId -> new rId so the caller can update
        XML references.
        """
        rid_mapping = {}

        for rId, rel in source_part.rels.items():
            if rel.reltype in _STRUCTURAL_REL_TYPES:
                continue

            if rel.is_external:
                new_rId = target_part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True,
                )
            elif rel.reltype == RT.IMAGE:
                image_blob = rel.target_part.blob
                image_stream = BytesIO(image_blob)
                _image_part, new_rId = target_part.get_or_add_image_part(image_stream)
            else:
                # For other internal rels (e.g. charts, media), copy the
                # blob as a generic Part.
                src_target = rel.target_part
                new_partname = target_package.next_partname(
                    _partname_to_template(src_target.partname),
                )
                new_part = Part(
                    new_partname,
                    src_target.content_type,
                    target_package,
                    blob=src_target.blob,
                )
                new_rId = target_part.relate_to(new_part, rel.reltype)

            if rId != new_rId:
                rid_mapping[rId] = new_rId

        return rid_mapping

    @staticmethod
    def _remap_rids(element, rid_mapping):
        """Walk an XML element tree and remap r:embed, r:link, r:id attributes."""
        r_attrs = (
            f"{{{_R_NS}}}embed",
            f"{{{_R_NS}}}link",
            f"{{{_R_NS}}}id",
        )
        for el in element.iter():
            for attr in r_attrs:
                val = el.get(attr)
                if val and val in rid_mapping:
                    el.set(attr, rid_mapping[val])

    # ------------------------------------------------------------------
    # ID generation
    # ------------------------------------------------------------------

    @staticmethod
    def _next_unique_id(prs: Presentation) -> int:
        """Return the next available id for sldMasterId / sldLayoutId elements.

        Both element types share the same id space.  Valid values start at
        2147483648 (0x80000000).  This helper scans all existing ids in
        ``sldMasterIdLst`` and every master's ``sldLayoutIdLst`` and returns
        ``max(existing) + 1``.
        """
        MIN_ID = 2147483648  # 0x80000000
        used: set[int] = set()

        prs_element = prs.part._element

        # Collect ids from sldMasterIdLst
        master_id_lst = prs_element.sldMasterIdLst
        if master_id_lst is not None:
            for entry in master_id_lst:
                val = entry.get("id")
                if val is not None:
                    used.add(int(val))

        # Collect ids from each master's sldLayoutIdLst
        for master in prs.slide_masters:
            layout_id_lst = master.part._element.sldLayoutIdLst
            if layout_id_lst is not None:
                for entry in layout_id_lst:
                    val = entry.get("id")
                    if val is not None:
                        used.add(int(val))

        return max(used | {MIN_ID - 1}) + 1

    # ------------------------------------------------------------------
    # Slide reordering
    # ------------------------------------------------------------------

    @staticmethod
    def _move_slide_to_index(prs: Presentation, target_index: int):
        """Move the last slide in the presentation to the given 0-based index.

        ``add_slide()`` always appends, so this helper is called afterwards to
        reposition the newly added slide.
        """
        sldIdLst = prs.slides._sldIdLst
        sldId_elements = list(sldIdLst)
        num_slides = len(sldId_elements)

        # Clamp index to valid range
        if target_index < 0:
            target_index = 0
        if target_index >= num_slides:
            # Already at the end, nothing to do
            return

        # The newly added slide is the last element
        new_sldId = sldId_elements[-1]
        sldIdLst.remove(new_sldId)

        if target_index == 0:
            sldIdLst.insert(0, new_sldId)
        else:
            # Insert after the element currently at target_index - 1
            ref_element = list(sldIdLst)[target_index - 1]
            ref_element.addnext(new_sldId)

    # ------------------------------------------------------------------
    # Existing helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _copy_slide_size(source_prs: Presentation, target_prs: Presentation):
        """Copy slide dimensions from source to target presentation."""
        try:
            target_prs.slide_width = source_prs.slide_width
            target_prs.slide_height = source_prs.slide_height
        except Exception:
            pass

    @staticmethod
    def _copy_images(source_slide: Slide, dest_slide: Slide):
        """Copy image parts and relationships from source slide to destination slide."""
        try:
            source_part = source_slide.part
            dest_part = dest_slide.part

            rId_mapping = {}

            for rel_id, rel in source_part.rels.items():
                if rel.reltype == RT.IMAGE:
                    image_part = rel.target_part
                    image_blob = image_part.blob
                    image_stream = BytesIO(image_blob)

                    result = dest_part.get_or_add_image_part(image_stream)

                    if isinstance(result, tuple):
                        new_image_part, new_rId = result
                        rId_mapping[rel_id] = new_rId
                        continue
                    else:
                        new_image_part = result

                    for new_rel_id, new_rel in dest_part.rels.items():
                        if (new_rel.reltype == RT.IMAGE and
                                new_rel.target_part == new_image_part):
                            rId_mapping[rel_id] = new_rel_id
                            break

            if rId_mapping:
                p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

                pics = dest_slide.element.findall(f'.//{{{p_ns}}}pic')
                for pic in pics:
                    blips = pic.findall(f'.//{{{a_ns}}}blip')
                    for blip in blips:
                        embed_attr = f'{{{_R_NS}}}embed'
                        old_rId = blip.get(embed_attr)

                        if old_rId and old_rId in rId_mapping:
                            blip.set(embed_attr, rId_mapping[old_rId])

        except Exception:
            pass


def _partname_to_template(partname):
    """Convert a PackURI like '/ppt/media/image3.png' to '/ppt/media/image%d.png'."""
    import re
    return re.sub(r'\d+(?=\.[^.]+$)', '%d', str(partname))

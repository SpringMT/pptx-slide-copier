"""Tests for SlideCopier."""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx_slide_copier import SlideCopier


@pytest.fixture
def sample_presentation():
    """Create a sample presentation with text and shapes."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add a slide with some text
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Add a text box
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Test slide content"

    # Set some formatting
    paragraph = text_frame.paragraphs[0]
    run = paragraph.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True

    return prs


@pytest.fixture
def temp_pptx_file(sample_presentation):
    """Save sample presentation to a temporary file."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        sample_presentation.save(tmp.name)
        yield tmp.name
        Path(tmp.name).unlink()


class TestSlideCopier:
    """Test cases for SlideCopier class."""

    def test_copy_slide_basic(self, sample_presentation):
        """Test basic slide copying."""
        source_prs = sample_presentation
        target_prs = Presentation()

        # Copy the first slide
        copied_slide = SlideCopier.copy_slide(source_prs, 0, target_prs)

        assert copied_slide is not None
        assert len(target_prs.slides) == 1

    def test_copy_slide_preserves_text(self, sample_presentation):
        """Test that text content is preserved."""
        source_prs = sample_presentation
        target_prs = Presentation()

        copied_slide = SlideCopier.copy_slide(source_prs, 0, target_prs)

        # Find text in copied slide
        found_text = False
        for shape in copied_slide.shapes:
            if hasattr(shape, "text") and "Test slide content" in shape.text:
                found_text = True
                break

        assert found_text, "Text content should be preserved"

    def test_copy_slide_preserves_size(self, sample_presentation):
        """Test that slide dimensions are preserved."""
        source_prs = sample_presentation
        target_prs = Presentation()

        SlideCopier.copy_slide(source_prs, 0, target_prs)

        assert target_prs.slide_width == source_prs.slide_width
        assert target_prs.slide_height == source_prs.slide_height

    def test_copy_multiple_slides(self, sample_presentation):
        """Test copying multiple slides."""
        source_prs = sample_presentation
        target_prs = Presentation()

        # Copy the same slide twice
        SlideCopier.copy_slide(source_prs, 0, target_prs)
        SlideCopier.copy_slide(source_prs, 0, target_prs)

        assert len(target_prs.slides) == 2

    def test_copy_slide_with_formatting(self, sample_presentation):
        """Test that text formatting is preserved."""
        source_prs = sample_presentation
        target_prs = Presentation()

        copied_slide = SlideCopier.copy_slide(source_prs, 0, target_prs)

        # Check that formatting is preserved
        for shape in copied_slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text == "Test slide content":
                            # Font properties should be preserved
                            assert run.font.bold
                            # Font name and size should be set
                            assert run.font.name is not None
                            assert run.font.size is not None

    def test_copy_slide_to_template_based_presentation(self, temp_pptx_file):
        """Test copying to a presentation based on the same template."""
        source_prs = Presentation(temp_pptx_file)

        # Create target from same template
        target_prs = Presentation(temp_pptx_file)

        # Remove all slides from target
        while len(target_prs.slides) > 0:
            rId = target_prs.slides._sldIdLst[0].rId
            target_prs.part.drop_rel(rId)
            del target_prs.slides._sldIdLst[0]

        assert len(target_prs.slides) == 0

        # Copy slide
        SlideCopier.copy_slide(source_prs, 0, target_prs)

        assert len(target_prs.slides) == 1

    def test_copy_slide_invalid_index(self, sample_presentation):
        """Test that invalid slide index raises appropriate error."""
        source_prs = sample_presentation
        target_prs = Presentation()

        with pytest.raises(IndexError):
            SlideCopier.copy_slide(source_prs, 999, target_prs)

    def test_copy_slide_preserves_shapes(self, sample_presentation):
        """Test that shapes are copied."""
        source_prs = sample_presentation

        target_prs = Presentation()
        copied_slide = SlideCopier.copy_slide(source_prs, 0, target_prs)

        # Should have same number of shapes (or close, due to layout elements)
        assert len(copied_slide.shapes) > 0


class TestLayoutPreservation:
    """Test that source slide layouts are faithfully copied."""

    @staticmethod
    def _make_source_with_custom_layout():
        """Create a source presentation whose layout name differs from the default template."""
        prs = Presentation()
        # The built-in default template has layouts like "Title Slide", "Title and Content", etc.
        # We'll use the first layout and record its name for later assertion.
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        textbox.text_frame.text = "hello"
        return prs

    def test_layout_name_preserved_across_different_templates(self):
        """Copying between presentations with different templates preserves layout name."""
        source_prs = self._make_source_with_custom_layout()
        source_layout_name = source_prs.slides[0].slide_layout.name

        # Target starts as a blank presentation (its own default template)
        target_prs = Presentation()

        SlideCopier.copy_slide(source_prs, 0, target_prs)

        copied_slide = target_prs.slides[0]
        assert copied_slide.slide_layout.name == source_layout_name

    def test_no_duplicate_parts_for_same_layout(self):
        """Copying two slides sharing the same layout should not duplicate master/layout parts."""
        source_prs = Presentation()
        layout = source_prs.slide_layouts[0]
        source_prs.slides.add_slide(layout)
        source_prs.slides.add_slide(layout)

        target_prs = Presentation()
        slides = SlideCopier.copy_slides(source_prs, target_prs, slide_indices=[0, 1])

        assert len(slides) == 2
        # Both copied slides must reference the same layout part
        assert slides[0].slide_layout.part is slides[1].slide_layout.part

    def test_layout_survives_save_and_reload(self):
        """Layout name is still correct after saving and reloading the file."""
        source_prs = self._make_source_with_custom_layout()
        source_layout_name = source_prs.slides[0].slide_layout.name

        target_prs = Presentation()
        SlideCopier.copy_slide(source_prs, 0, target_prs)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            target_prs.save(tmp.name)
            tmp_path = tmp.name

        try:
            reloaded = Presentation(tmp_path)
            assert len(reloaded.slides) == 1
            assert reloaded.slides[0].slide_layout.name == source_layout_name
        finally:
            Path(tmp_path).unlink()

    def test_copy_slides_convenience_method(self):
        """copy_slides copies all slides when slide_indices is None."""
        source_prs = Presentation()
        source_prs.slides.add_slide(source_prs.slide_layouts[0])
        source_prs.slides.add_slide(source_prs.slide_layouts[1])

        target_prs = Presentation()
        slides = SlideCopier.copy_slides(source_prs, target_prs)

        assert len(slides) == 2
        assert len(target_prs.slides) == 2

    def test_copy_layouts_then_slides(self):
        """copy_layouts() followed by copy_slide() with layout_map works correctly."""
        source_prs = Presentation()
        layout_0 = source_prs.slide_layouts[0]
        layout_1 = source_prs.slide_layouts[1]
        source_prs.slides.add_slide(layout_0)
        source_prs.slides.add_slide(layout_1)

        target_prs = Presentation()
        layout_map = SlideCopier.copy_layouts(source_prs, target_prs)

        # All source layout names should be in the map
        for master in source_prs.slide_masters:
            for layout in master.slide_layouts:
                assert layout.name in layout_map

        # Copy slides using the layout_map
        s0 = SlideCopier.copy_slide(source_prs, 0, target_prs, _layout_map=layout_map)
        s1 = SlideCopier.copy_slide(source_prs, 1, target_prs, _layout_map=layout_map)

        assert len(target_prs.slides) == 2
        assert s0.slide_layout.name == layout_0.name
        assert s1.slide_layout.name == layout_1.name

    def test_target_has_two_themes(self):
        """After copy_layouts with different themes, target should have original + source theme."""
        from pptx.opc.constants import CONTENT_TYPE as CT
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        def _count_theme_parts(prs):
            parts = set()
            for rel in prs.part.package.iter_rels():
                try:
                    if rel.target_part.content_type == CT.OFC_THEME:
                        parts.add(rel.target_part.partname)
                except Exception:
                    pass
            return len(parts)

        source_prs = Presentation()
        source_prs.slides.add_slide(source_prs.slide_layouts[0])

        # ソースのテーマ blob を変更して異なるテーマにする
        source_theme_part = source_prs.slide_masters[0].part.part_related_by(RT.THEME)
        source_theme_part._blob = source_theme_part.blob + b"<!-- modified -->"

        target_prs = Presentation()
        original_theme_count = _count_theme_parts(target_prs)

        SlideCopier.copy_layouts(source_prs, target_prs)

        total_theme_count = _count_theme_parts(target_prs)
        assert total_theme_count == original_theme_count + 1

    def test_same_theme_no_duplicate(self):
        """同じテーマを持つソースとターゲットでcopy_layouts後にテーマ数が増えないこと。"""
        from pptx.opc.constants import CONTENT_TYPE as CT

        def _count_theme_parts(prs):
            parts = set()
            for rel in prs.part.package.iter_rels():
                try:
                    if rel.target_part.content_type == CT.OFC_THEME:
                        parts.add(rel.target_part.partname)
                except Exception:
                    pass
            return len(parts)

        # 同じデフォルトテンプレートから作成
        source_prs = Presentation()
        source_prs.slides.add_slide(source_prs.slide_layouts[0])

        target_prs = Presentation()
        original_theme_count = _count_theme_parts(target_prs)

        SlideCopier.copy_layouts(source_prs, target_prs)

        total_theme_count = _count_theme_parts(target_prs)
        assert total_theme_count == original_theme_count

    def test_copy_slide_target_index_insert_at_beginning(self):
        """target_slide_index=0 inserts the copied slide at the beginning."""
        target_prs = Presentation()
        layout = target_prs.slide_layouts[0]
        # Add two existing slides with identifiable text
        s1 = target_prs.slides.add_slide(layout)
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "existing-1"
        s2 = target_prs.slides.add_slide(layout)
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "existing-2"

        source_prs = Presentation()
        src_slide = source_prs.slides.add_slide(source_prs.slide_layouts[0])
        src_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "copied"

        SlideCopier.copy_slide(source_prs, 0, target_prs, target_slide_index=0)

        assert len(target_prs.slides) == 3
        # The copied slide should be first
        texts = []
        for slide in target_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
                    break
        assert texts == ["copied", "existing-1", "existing-2"]

    def test_copy_slide_target_index_insert_in_middle(self):
        """target_slide_index=1 inserts the copied slide at position 1."""
        target_prs = Presentation()
        layout = target_prs.slide_layouts[0]
        s1 = target_prs.slides.add_slide(layout)
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "existing-1"
        s2 = target_prs.slides.add_slide(layout)
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "existing-2"

        source_prs = Presentation()
        src_slide = source_prs.slides.add_slide(source_prs.slide_layouts[0])
        src_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "copied"

        SlideCopier.copy_slide(source_prs, 0, target_prs, target_slide_index=1)

        assert len(target_prs.slides) == 3
        texts = []
        for slide in target_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
                    break
        assert texts == ["existing-1", "copied", "existing-2"]

    def test_copy_slide_target_index_none_appends(self):
        """target_slide_index=None (default) appends at the end."""
        target_prs = Presentation()
        layout = target_prs.slide_layouts[0]
        s1 = target_prs.slides.add_slide(layout)
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "existing-1"

        source_prs = Presentation()
        src_slide = source_prs.slides.add_slide(source_prs.slide_layouts[0])
        src_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "copied"

        SlideCopier.copy_slide(source_prs, 0, target_prs)

        assert len(target_prs.slides) == 2
        texts = []
        for slide in target_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
                    break
        assert texts == ["existing-1", "copied"]

    def test_copy_slides_target_index(self):
        """copy_slides with target_slide_index inserts slides sequentially."""
        target_prs = Presentation()
        layout = target_prs.slide_layouts[0]
        s1 = target_prs.slides.add_slide(layout)
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "existing-1"
        s2 = target_prs.slides.add_slide(layout)
        s2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "existing-2"

        source_prs = Presentation()
        src1 = source_prs.slides.add_slide(source_prs.slide_layouts[0])
        src1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "copied-A"
        src2 = source_prs.slides.add_slide(source_prs.slide_layouts[0])
        src2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "copied-B"

        SlideCopier.copy_slides(source_prs, target_prs, target_slide_index=1)

        assert len(target_prs.slides) == 4
        texts = []
        for slide in target_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
                    break
        assert texts == ["existing-1", "copied-A", "copied-B", "existing-2"]

    def test_same_theme_layout_map_uses_existing(self):
        """テーマ同一時、layout_mapの値がターゲットの既存レイアウトを指すこと。"""
        source_prs = Presentation()
        source_prs.slides.add_slide(source_prs.slide_layouts[0])

        target_prs = Presentation()
        # ターゲットの既存レイアウトを記録
        existing_layouts = {sl.name: sl for master in target_prs.slide_masters
                           for sl in master.slide_layouts}

        layout_map = SlideCopier.copy_layouts(source_prs, target_prs)

        # layout_mapの各値がターゲットの既存レイアウトと同一であること
        for name, layout in layout_map.items():
            if name in existing_layouts:
                assert layout is existing_layouts[name]


class TestMasterLayoutIdAttributes:
    """Test that copied sldMasterId and sldLayoutId elements have valid id attributes."""

    @staticmethod
    def _make_different_theme_source():
        """Create a source presentation with a modified theme blob."""
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        prs = Presentation()
        theme_part = prs.slide_masters[0].part.part_related_by(RT.THEME)
        theme_part._blob = theme_part.blob + b"<!-- different -->"
        prs.slides.add_slide(prs.slide_layouts[0])
        return prs

    def test_copy_slide_different_theme_has_master_id(self):
        """sldMasterId elements must have an id attribute after copy with different theme."""
        source_prs = self._make_different_theme_source()
        target_prs = Presentation()

        SlideCopier.copy_slide(source_prs, 0, target_prs)

        master_id_lst = target_prs.part._element.sldMasterIdLst
        for entry in master_id_lst:
            assert entry.get("id") is not None, "sldMasterId missing id attribute"

    def test_copy_slide_different_theme_has_layout_id(self):
        """sldLayoutId elements must have an id attribute after copy with different theme."""
        source_prs = self._make_different_theme_source()
        target_prs = Presentation()

        SlideCopier.copy_slide(source_prs, 0, target_prs)

        for master in target_prs.slide_masters:
            layout_id_lst = master.part._element.sldLayoutIdLst
            if layout_id_lst is not None:
                for entry in layout_id_lst:
                    assert entry.get("id") is not None, "sldLayoutId missing id attribute"

    def test_copy_slides_different_theme_has_master_id(self):
        """sldMasterId elements must have an id attribute after copy_slides with different theme."""
        source_prs = self._make_different_theme_source()
        target_prs = Presentation()

        SlideCopier.copy_slides(source_prs, target_prs)

        master_id_lst = target_prs.part._element.sldMasterIdLst
        for entry in master_id_lst:
            assert entry.get("id") is not None, "sldMasterId missing id attribute"

    def test_ids_are_unique(self):
        """All sldMasterId and sldLayoutId id values must be unique across the presentation."""
        source_prs = self._make_different_theme_source()
        target_prs = Presentation()

        SlideCopier.copy_slide(source_prs, 0, target_prs)

        all_ids = []
        master_id_lst = target_prs.part._element.sldMasterIdLst
        if master_id_lst is not None:
            for entry in master_id_lst:
                val = entry.get("id")
                if val is not None:
                    all_ids.append(int(val))

        for master in target_prs.slide_masters:
            layout_id_lst = master.part._element.sldLayoutIdLst
            if layout_id_lst is not None:
                for entry in layout_id_lst:
                    val = entry.get("id")
                    if val is not None:
                        all_ids.append(int(val))

        assert len(all_ids) == len(set(all_ids)), f"Duplicate ids found: {all_ids}"

    def test_copy_slide_different_theme_with_target_index_has_ids(self):
        """File must not be corrupted when using target_slide_index with different theme."""
        source_prs = self._make_different_theme_source()
        target_prs = Presentation()
        target_prs.slides.add_slide(target_prs.slide_layouts[0])
        target_prs.slides.add_slide(target_prs.slide_layouts[0])

        SlideCopier.copy_slide(source_prs, 0, target_prs, target_slide_index=1)

        master_id_lst = target_prs.part._element.sldMasterIdLst
        for entry in master_id_lst:
            assert entry.get("id") is not None, "sldMasterId missing id attribute"

        # Verify file can be saved and reloaded
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            target_prs.save(tmp.name)
            tmp_path = tmp.name

        try:
            reloaded = Presentation(tmp_path)
            assert len(reloaded.slides) == 3
        finally:
            Path(tmp_path).unlink()
